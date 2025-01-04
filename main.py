import os
from pptx import Presentation
import re
import tkinter as tk
from tkinter import filedialog, messagebox

# 读取设置
try:
    from setting import FULL_MATCH, IGNORE_CASE
except ImportError:
    print("未找到 setting.py 文件，将使用默认设置。")


def search_in_ppt(folder_path, search_term):
    """
    搜索指定文件夹下所有 PPT 文件中的指定字符，并返回匹配的文件名和定位页。

    :param folder_path: 文件夹路径
    :param search_term: 要搜索的字符
    :return: list，格式为 ["文件名 - 页码"]
    """
    result = []

    # 正则表达式设置
    flags = re.IGNORECASE if IGNORE_CASE else 0
    if FULL_MATCH:
        pattern = re.compile(rf"\\b{re.escape(search_term)}\\b", flags)
    else:
        pattern = re.compile(re.escape(search_term), flags)

    # 遍历文件夹中的 PPT 文件
    for file_name in os.listdir(folder_path):
        if file_name.endswith(".pptx"):
            file_path = os.path.join(folder_path, file_name)
            try:
                ppt = Presentation(file_path)
                for i, slide in enumerate(ppt.slides):
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            text = shape.text
                            if pattern.search(text):
                                result.append(f"{file_name} - 第{i + 1}页")
            except Exception as e:
                print(f"处理文件 {file_name} 时出错: {e}")

    return result


def browse_folder():
    folder = filedialog.askdirectory()
    folder_path_var.set(folder)


def search():
    folder = folder_path_var.get()
    term = search_term_var.get()

    if not folder or not term:
        messagebox.showwarning("警告", "请填写完整的文件夹路径和搜索字符！")
        return

    matches = search_in_ppt(folder, term)

    if matches:
        result_text.delete(1.0, tk.END)
        result_text.insert(tk.END, "\n".join(matches))
    else:
        result_text.delete(1.0, tk.END)
        result_text.insert(tk.END, "未找到匹配项。")


# 创建主窗口
root = tk.Tk()
root.title("PPT 搜索工具")
root.geometry("600x400")

# 文件夹选择
folder_path_var = tk.StringVar()
search_term_var = tk.StringVar()

frame = tk.Frame(root)
frame.pack(pady=10, padx=10, fill=tk.X)

folder_label = tk.Label(frame, text="文件夹路径:")
folder_label.pack(side=tk.LEFT)

folder_entry = tk.Entry(frame, textvariable=folder_path_var, width=50)
folder_entry.pack(side=tk.LEFT, padx=5)

browse_button = tk.Button(frame, text="浏览", command=browse_folder)
browse_button.pack(side=tk.LEFT)

# 搜索字符输入
term_label = tk.Label(root, text="搜索字符:")
term_label.pack()

term_entry = tk.Entry(root, textvariable=search_term_var, width=50)
term_entry.pack(pady=5)

# 搜索按钮
search_button = tk.Button(root, text="搜索", command=search)
search_button.pack(pady=5)

# 结果显示框
result_text = tk.Text(root, height=15, wrap=tk.WORD)
result_text.pack(pady=10, padx=10, fill=tk.BOTH, expand=True)

# 运行主循环
root.mainloop()
