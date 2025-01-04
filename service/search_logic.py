import os
from pptx import Presentation
import re

# 读取设置
try:
    from setting.setting import FULL_MATCH, IGNORE_CASE
except ImportError:
    print("未找到 setting.py 文件，将使用默认设置。")


def search_in_ppt(file_or_folder_path, search_term, is_folder=True):
    """
    搜索指定文件夹或单个 PPT 文件中的指定字符，并返回匹配的文件名和定位页，以及匹配内容。

    :param file_or_folder_path: 文件夹路径或单个文件路径
    :param search_term: 要搜索的字符
    :param is_folder: 是否为文件夹路径，默认为 True
    :return: list，格式为 ["文件名 - 页码 - 匹配内容"]
    """
    result = []

    # 正则表达式设置
    flags = re.IGNORECASE if IGNORE_CASE else 0
    if FULL_MATCH:
        pattern = re.compile(rf"\\b{re.escape(search_term)}\\b", flags)
    else:
        pattern = re.compile(re.escape(search_term), flags)

    def process_ppt(file_path, file_name):
        try:
            ppt = Presentation(file_path)
            for i, slide in enumerate(ppt.slides):
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text = shape.text
                        # 查找匹配
                        matches = pattern.finditer(text)
                        for match in matches:
                            start_idx = match.start()
                            # 提取匹配行的上下文，最多包含8个字符
                            line_start = max(start_idx - 8, 0)
                            line_end = min(start_idx + 8 + len(search_term), len(text))
                            matching_text = text[line_start:line_end]

                            # 记录匹配信息：文件名 - 页码 - 匹配内容
                            result.append(f"{file_name} - 第{i + 1}页 - \"{matching_text}\"")
        except Exception as e:
            print(f"处理文件 {file_name} 时出错: {e}")

    if is_folder:
        for file_name in os.listdir(file_or_folder_path):
            if file_name.endswith(".pptx"):
                file_path = os.path.join(file_or_folder_path, file_name)
                process_ppt(file_path, file_name)
    else:
        file_name = os.path.basename(file_or_folder_path)
        process_ppt(file_or_folder_path, file_name)

    return result
